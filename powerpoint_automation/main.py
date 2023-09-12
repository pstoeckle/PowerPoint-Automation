# SPDX-FileCopyrightText: 2022 Patrick Stöckle.
# SPDX-License-Identifier: Apache-2.0
"""
Convert.
"""
from logging import INFO, basicConfig, getLogger
from os.path import isfile
from pathlib import Path
from sys import platform, stdout
from typing import List, Optional

from pptx import Presentation
from pptx.shapes.autoshape import Shape

from powerpoint_automation import __version__
from powerpoint_automation.logic.add_git_info import add_git_info_internal
from powerpoint_automation.logic.add_meta_data import add_meta_data_internal
from powerpoint_automation.logic.convert_presentations import (
    convert_presentations_internal,
)
from powerpoint_automation.logic.create_txt_for_powerpoint import process_pptx_file
from powerpoint_automation.logic.remove_picture import remove_picture_internal
from typer import Argument, Exit, Option, Typer, echo, style
from typer.colors import RED

LINUX_LIBREOFFICE = "/usr/bin/libreoffice"
MAC_OS_SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
LIBRE_OFFICE = "NOT_SET"

_LOGGER = getLogger(__name__)
basicConfig(
    format="%(levelname)s: %(asctime)s: %(name)s: %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
    level=INFO,
    stream=stdout,
)

_INPUT_DIRECTORY_OPTION = Argument(
    ".", exists=True, resolve_path=True
)

def _version_callback(value: bool) -> None:
    if value:
        echo(f"powerpoint-automation {__version__}")
        raise Exit()


app = Typer()


def error_echo(s: str) -> None:
    """

    :param s:
    :return:
    """
    echo(style(s, fg=RED), err=True)


@app.callback()
def _call_back(
    _: bool = Option(
        None,
        "--version",
        is_flag=True,
        callback=_version_callback,
        expose_value=False,
        is_eager=True,
        help="Version",
    )
) -> None:
    """
    Scripts to automate some task in the context of PowerPoint.
    """


def _set_libreoffice() -> None:
    global LIBRE_OFFICE

    if (platform == "linux" or platform == "linux2") and isfile(LINUX_LIBREOFFICE):
        LIBRE_OFFICE = LINUX_LIBREOFFICE
    elif platform == "darwin" and isfile(MAC_OS_SOFFICE):
        LIBRE_OFFICE = MAC_OS_SOFFICE
    else:
        _LOGGER.debug(f"Using PowerPoint for Windows.")


_set_libreoffice()


@app.command()
def convert_presentations(
    input_directory_path: Path = _INPUT_DIRECTORY_OPTION,
    output_directory_path: Path = Option(
        "dist",
        "--output-directory",
        "-o",
        file_okay=False,
        resolve_path=True,
    ),
    libre_office: Path = Option(
        LIBRE_OFFICE,
        "--libre-office",
        "-L",
        resolve_path=True,
        dir_okay=False,
    ),
    skip_file: Optional[List[str]] = Option(None, "--skip-file", "-s"),
) -> None:
    """
    Converts PowerPoint files to PDFs.
    This is especially handy if you want to create many PDFs without exporting them manually one by one.
    """
    skip_files = frozenset() if skip_file is None else frozenset(s for s in skip_file)
    convert_presentations_internal(
        input_directory_path, libre_office, output_directory_path, skip_files
    )


@app.command()
def replace_date(
    input_directory_path: Path = _INPUT_DIRECTORY_OPTION,
    old_year: int = Option(2020, "--old-year", "-O"),
    new_year: int = Option(2021, "--new-year", "-N"),
) -> None:
    """
    Replace a date in the slides, e.g., 2020 -> 2021.
    This is especially handy at the beginning of a new semester where you have to update the slides of the
    previous year.
    """
    pptx_files = sorted([
        f
        for f in input_directory_path.iterdir()
        if f.is_file() and f.suffix.casefold() == ".pptx" and not f.stem.startswith("~")
    ])
    for input_file in pptx_files:
        pres = Presentation(input_file)
        rewrite_file = False
        for shape in pres.slide_master.shapes:
            rewrite_file = _replace_old_year_in_shape(shape, new_year, old_year) or rewrite_file
        for layout in pres.slide_master.slide_layouts:
            for shape in layout.shapes:
                rewrite_file = _replace_old_year_in_shape(shape, new_year, old_year) or rewrite_file
        for slide_no, slide in enumerate(pres.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                old_text = shape.text.casefold().replace(" ", "")
                if (
                    "|securityengineering|" in old_text
                    and f"summer{old_year}" in old_text
                ):
                    shape.text = f"Prof. Dr. Alexander Pretschner (I4) | Security Engineering | Summer {new_year}"
                    rewrite_file = True
                    _LOGGER.debug(f"{input_file}: Changing field on slide  {slide_no}")
                    continue
                if (
                    "in2178–securityengineering" in old_text
                    and f"summersemester{old_year}" in old_text
                ):
                    shape.text = shape.text.replace(str(old_year), str(new_year))
                    rewrite_file = True
                    _LOGGER.debug(f"{input_file}: Changing field on slide  {slide_no}")
                    continue

        if rewrite_file:
            echo(f"Rewrite file {input_file}")
            pres.save(input_file)


def _replace_old_year_in_shape(shape: Shape, new_year: int, old_year: int )-> bool:
    if shape.has_text_frame and str(old_year) in shape.text:
        shape.text = shape.text.replace(str(old_year), str(new_year))
        return True
    return False



@app.command()
def add_git_metadata_as_footer(input_directory_path: Path = _INPUT_DIRECTORY_OPTION) -> None:
    """
    Adds a footer to the PowerPoint slides with the latest commit's hash and date.
    """
    add_git_info_internal(input_directory_path)


@app.command()
def add_git_metadata_as_presentation_metadata(
    input_directory_path: Path = _INPUT_DIRECTORY_OPTION,
    author: Optional[List[str]] = Option(None, "--author", "-a"),
) -> None:
    """
    This commands adds git metadata to the PowerPoint slides, e.g., the commit, the authors, etc.
    """
    if author is None:
        author = []
    add_meta_data_internal(author, input_directory_path)


@app.command()
def remove_picture(
    input_directory_path: Path = _INPUT_DIRECTORY_OPTION,
    hash_value: Optional[List[str]] = Option(None, "--hash-value", "-S"),
    inplace: bool = Option(False, "--inplace", "-i", is_flag=True),
) -> None:
    """
    Remove the pictures from all slides.
    """
    if hash_value is None or len(hash_value) == 0:
        error_echo("No hashes ...")
        raise Exit(1)
    hashes_set = frozenset(h.casefold() for h in hash_value)
    remove_picture_internal(hashes_set, inplace, input_directory_path)


@app.command()
def git_pptx_diff(
    filename: Path = Argument(
        "filename", exists=True, resolve_path=True, dir_okay=False
    )
) -> None:
    """

    :return:
    """
    echo(process_pptx_file(filename))


if __name__ == "__main__":
    app()
