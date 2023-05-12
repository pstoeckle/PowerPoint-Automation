# SPDX-FileCopyrightText: 2022 Patrick Stöckle.
# SPDX-License-Identifier: Apache-2.0
"""
Git info.
"""
from logging import getLogger
from pathlib import Path
from subprocess import check_output

from pptx import Presentation
from pptx.util import Cm, Pt

_LOGGER = getLogger(__name__)


def add_git_info_internal(input_directory_path: Path) -> None:
    """

    :param input_directory_path:
    :return:
    """
    pptx_files = [
        f
        for f in input_directory_path.iterdir()
        if f.is_file() and f.suffix.casefold() == ".pptx" and not f.stem.startswith("~")
    ]
    for input_file in pptx_files:
        pres = Presentation(input_file)
        _LOGGER.info(f"Processing file {input_file}")
        commit_sha = check_output(
            ["git", "log", "-n", "1", "--pretty=format:%h", "--", input_file]
        ).decode("utf8")
        commit_date = check_output(
            ["git", "log", "-n", "1", "--pretty=format:%aI", "--", input_file]
        ).decode("utf8")
        for slide_no, slide in enumerate(pres.slides):
            text_box = slide.shapes.add_textbox(Cm(24), Cm(17.33), Cm(7), Cm(1))
            tf = text_box.text_frame
            p = tf.add_paragraph()
            p.font.size = Pt(10)
            p.text = f"{commit_date} | {commit_sha}"

        pres.save(input_file)
        _LOGGER.info(f"{input_file}: Done")
