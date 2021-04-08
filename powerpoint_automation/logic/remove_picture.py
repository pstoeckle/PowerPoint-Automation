"""
Remove pictures from PPTX.
"""
from logging import getLogger
from pathlib import Path
from typing import AbstractSet

from pptx import Presentation
from pptx.shapes.picture import Picture

_LOGGER = getLogger(__name__)


def remove_picture_internal(
    hashes_set: AbstractSet[str], inplace: bool, input_directory_path: Path
) -> None:
    pptx_files = [
        f
        for f in input_directory_path.iterdir()
        if f.is_file() and f.suffix.casefold() == ".pptx" and not f.stem.startswith("~")
    ]
    for input_file in pptx_files:
        pres = Presentation(input_file)
        rewrite_file = False
        for slide in pres.slides:
            shapes_to_delete = []
            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    if shape.image.sha1.casefold() in hashes_set:
                        shapes_to_delete.append(shape)
            if len(shapes_to_delete) > 0:
                rewrite_file = True
                for shape_to_delete in shapes_to_delete:
                    old_pic = shape_to_delete._element
                    old_pic.getparent().remove(old_pic)
        if rewrite_file:
            if inplace:
                _LOGGER.info(f"Rewrite file {input_file}")
                pres.save(input_file)
            else:
                new_file_name = (
                    str(input_file)
                    .replace(".pptx", ".out.pptx")
                    .replace(".PPTX", ".out.pptx")
                )
                _LOGGER.info(f"Write file {new_file_name}")
                pres.save(new_file_name)
