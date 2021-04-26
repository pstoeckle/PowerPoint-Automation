#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Commands.
"""

from os import linesep, listdir, mkdir, path

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.text.text import TextFrame

END = "END"

START = "START"

_INDENT = "    "
_TEXTS = "texts"
_PPTX = ".pptx"


def process_all_pptx_files_in_folder(folder: str) -> None:
    """

    :return:
    """
    if not path.isdir(_TEXTS):
        mkdir(_TEXTS)
    pptx_files = [
        path.join(folder, f)
        for f in listdir(folder)
        if f.endswith(_PPTX) and not f.startswith("~")
    ]
    for pptx_filename in pptx_files:
        process_pptx_file(pptx_filename, True)


def process_pptx_file(pptx_filename: str, print_to_file: bool = False) -> str:
    """

    :param pptx_filename:
    :param print_to_file:
    :return:
    """
    string_to_write = ""
    prs = Presentation(pptx_filename)
    for i, slide in enumerate(prs.slides):
        string_to_write += f"{START}: slide {i}" + linesep
        for j, shape in enumerate(slide.shapes):
            string_to_write += _handle_shape(j, shape, 1)
        string_to_write += f"{END}: slide {i}" + linesep + linesep
    pptx_filename_replace = path.join(_TEXTS, pptx_filename.replace(_PPTX, ".txt"))
    if print_to_file:
        with open(pptx_filename_replace, "w") as f_write:
            f_write.write(string_to_write)
    return string_to_write


def _handle_shape(number_of_shape: int, shape: Shape, number_of_indent: int) -> str:
    """

    :param number_of_shape:
    :param shape:
    :param number_of_indent:
    :return:
    """
    indent = _get_indent(number_of_indent)
    indent_1 = _get_indent(number_of_indent + 1)
    string_to_write = indent + "{}: shape {}".format(START, number_of_shape) + linesep
    if shape.has_text_frame:
        text_frame: TextFrame = shape.text_frame
        for k, paragraph in enumerate(text_frame.paragraphs):
            t = paragraph.text.replace("\n", linesep + indent_1)
            if paragraph.font.bold and t != "":
                t = "**{}**".format(t)
            if paragraph.font.italic and t != "":
                t = "_{}_".format(t)
            string_to_write += indent_1 + t + linesep

    if isinstance(shape, GroupShape):
        for i, sub_shape in enumerate(shape.shapes):
            string_to_write += _handle_shape(i, sub_shape, number_of_indent + 1)
    string_to_write += indent + f"{END}: shape {number_of_shape}" + linesep
    return string_to_write


def _get_indent(n: int) -> str:
    """

    :param n:
    :return:
    """
    return "".join(_INDENT for _ in range(0, n))
